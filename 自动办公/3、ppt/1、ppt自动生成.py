from pptx import Presentation
from pptx.util import Inches,Pt
#安装使用pip install pptx

ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[1]) #在ppt中插入一个幻灯片

body_shape = slide.shapes.placeholders
# body_shape[0].text = '这是占位符[0]'
# body_shape[1].text = '这是占位符[1]'

title_shape = slide.shapes.title
title_shape.text = '这是标题'
# subtitle = slide.shapes.placeholders[1]   #取出本页第二个文本框
# subtitle.text = '这是文本框‘ #在第二个文本框中写入文字

new_paragraph = body_shape[1].text_frame.add_paragraph()
new_paragraph.text = '新段落'
new_paragraph.font.bold = True      #字体加粗
new_paragraph.font.italic = True    #文字斜体
new_paragraph.font.size = Pt(15)    #文字大小
new_paragraph.font.underline = True #文字下划线

#添加新的文本框
left = Inches(2)
top = Inches(2)
width = Inches(3)
height = Inches(3)

textbox = slide.shapes.add_textbox(left, top, width, height)
textbox.text = '这是新文本框'
new_para = textbox.text_frame.add_paragraph()
new_para.text = '这是新文本框里的第二段'

#插入图片
left = Inches(2)
top = Inches(2)
width = Inches(3)
height = Inches(3)

pic = slide.shapes.add_picture('b.jpg', left, top, width, height)

#插入表格
rows = 2
cols = 2
left = Inches(1)
top = Inches(1)
width = Inches(4)
height = Inches(4)
table = slide.shapes.add_table(rows, cols, left, width, height)
table.columns[0].width = Inches(1)
table.columns[1].width = Inches(3)
table.cell(0, 0).text = '1'
table.cell(0, 1).text = '2'
table.cell(1, 0).text = '3'
table.cell(1, 1).text = '4'

ppt.save('test.pptx')
